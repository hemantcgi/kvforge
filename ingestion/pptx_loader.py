"""PowerPoint document loader that extracts slides and speaker notes.

Depends on ``python-pptx`` for text extraction. Install with ``pip install python-pptx``.
"""
import hashlib
from pathlib import Path
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


def _section_hash(text: str) -> str:
    """Compute SHA-256 hash of text."""
    return hashlib.sha256(text.encode("utf-8")).hexdigest()


class PptxLoader:
    """Load a PowerPoint file and extract slides and speaker notes as chunks.

    Each slide produces one chunk for the body text (title + content shapes).
    Speaker notes (if present) produce a companion chunk with is_speaker_notes=True.
    Both chunks from the same slide share the same section_hash.

    Args:
        None
    """

    def __init__(self):
        pass

    def load(self, source: str) -> list[dict]:
        """Read a PowerPoint file and return slides and notes as chunk dicts.

        Args:
            source: Path to the PPTX file.

        Returns:
            List of dicts with the shape::

                {
                    "text": str,
                    "metadata": {
                        "slide_number": int,      # 1-indexed slide number
                        "slide_title": str,       # title text from the slide
                        "is_speaker_notes": bool, # True if this chunk is speaker notes
                        "section_hash": str,      # SHA-256 hash (shared by body + notes from same slide)
                        "source": str,            # filename (not full path)
                        "chunk_id": int           # global 0-indexed chunk counter
                    }
                }
        """
        path = Path(source)
        prs = Presentation(str(path))
        chunks: list[dict] = []
        chunk_id = 0

        for slide_num, slide in enumerate(prs.slides, start=1):
            # Extract title from title shapes
            slide_title = ""
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.has_text_frame:
                    # Check if it's a title shape using placeholder type
                    if hasattr(shape, "placeholder_format"):
                        try:
                            from pptx.enum.shapes import PP_PLACEHOLDER
                            if shape.placeholder_format.type in (PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE):
                                slide_title = shape.text.strip()
                                break
                        except (AttributeError, ImportError):
                            pass

            # Extract body text from non-title shapes (skip pictures)
            body_parts = []
            for shape in slide.shapes:
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    continue
                if hasattr(shape, "text") and shape.has_text_frame:
                    # Check if it's a title shape
                    if hasattr(shape, "placeholder_format"):
                        try:
                            from pptx.enum.shapes import PP_PLACEHOLDER
                            if shape.placeholder_format.type in (PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE):
                                continue
                        except (AttributeError, ImportError):
                            pass
                    text = shape.text.strip()
                    if text:
                        body_parts.append(text)

            body_text = " ".join(body_parts)

            # Extract speaker notes
            notes_text = ""
            if slide.has_notes_slide:
                notes_frame = slide.notes_slide.notes_text_frame
                notes_text = notes_frame.text.strip() if notes_frame else ""

            # Compute section_hash over combined body + notes
            combined_text = f"{body_text} {notes_text}".strip()
            section_hash = _section_hash(combined_text)

            # Add body chunk if there's body text
            if body_text:
                chunks.append({
                    "text": body_text,
                    "metadata": {
                        "slide_number": slide_num,
                        "slide_title": slide_title,
                        "is_speaker_notes": False,
                        "section_hash": section_hash,
                        "source": path.name,
                        "chunk_id": chunk_id,
                    },
                })
                chunk_id += 1

            # Add notes chunk if there are speaker notes
            if notes_text:
                chunks.append({
                    "text": notes_text,
                    "metadata": {
                        "slide_number": slide_num,
                        "slide_title": slide_title,
                        "is_speaker_notes": True,
                        "section_hash": section_hash,
                        "source": path.name,
                        "chunk_id": chunk_id,
                    },
                })
                chunk_id += 1

        return chunks
