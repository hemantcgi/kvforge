import sys
from pathlib import Path
sys.path.insert(0, str(Path(__file__).parent.parent))

import pytest
from pptx import Presentation
from pptx.util import Inches, Pt


def _make_pptx(tmp_path, slides: list[dict]) -> str:
    """slides: list of {title, body, notes}"""
    prs = Presentation()
    layout = prs.slide_layouts[1]  # title + content
    for s in slides:
        slide = prs.slides.add_slide(layout)
        slide.shapes.title.text = s.get("title", "")
        slide.placeholders[1].text = s.get("body", "")
        if s.get("notes"):
            slide.notes_slide.notes_text_frame.text = s["notes"]
    path = tmp_path / "test.pptx"
    prs.save(str(path))
    return str(path)


def test_pptx_slide_number(tmp_path):
    path = _make_pptx(tmp_path, [
        {"title": "Slide One", "body": "Content A"},
        {"title": "Slide Two", "body": "Content B"},
    ])
    from ingestion.pptx_loader import PptxLoader
    chunks = PptxLoader().load(path)
    slide_numbers = [c["metadata"]["slide_number"] for c in chunks if not c["metadata"]["is_speaker_notes"]]
    assert 1 in slide_numbers
    assert 2 in slide_numbers


def test_pptx_slide_title_in_metadata(tmp_path):
    path = _make_pptx(tmp_path, [{"title": "My Title", "body": "Body text."}])
    from ingestion.pptx_loader import PptxLoader
    chunks = PptxLoader().load(path)
    main_chunk = next(c for c in chunks if not c["metadata"]["is_speaker_notes"])
    assert main_chunk["metadata"]["slide_title"] == "My Title"


def test_pptx_body_text_in_chunk(tmp_path):
    path = _make_pptx(tmp_path, [{"title": "T", "body": "Important content here."}])
    from ingestion.pptx_loader import PptxLoader
    chunks = PptxLoader().load(path)
    texts = " ".join(c["text"] for c in chunks)
    assert "Important content here." in texts


def test_pptx_speaker_notes_as_companion_chunk(tmp_path):
    path = _make_pptx(tmp_path, [{"title": "T", "body": "B", "notes": "These are the notes."}])
    from ingestion.pptx_loader import PptxLoader
    chunks = PptxLoader().load(path)
    notes_chunks = [c for c in chunks if c["metadata"]["is_speaker_notes"]]
    assert len(notes_chunks) == 1
    assert "These are the notes." in notes_chunks[0]["text"]
    assert notes_chunks[0]["metadata"]["slide_number"] == 1


def test_pptx_section_hash_present(tmp_path):
    path = _make_pptx(tmp_path, [{"title": "T", "body": "B"}])
    from ingestion.pptx_loader import PptxLoader
    chunks = PptxLoader().load(path)
    for c in chunks:
        assert "section_hash" in c["metadata"]
        assert len(c["metadata"]["section_hash"]) == 64


def test_pptx_source_file_in_metadata(tmp_path):
    path = _make_pptx(tmp_path, [{"title": "T", "body": "B"}])
    from ingestion.pptx_loader import PptxLoader
    chunks = PptxLoader().load(path)
    assert all(c["metadata"]["source"] == Path(path).name for c in chunks)


def test_pptx_empty_presentation(tmp_path):
    prs = Presentation()
    path = tmp_path / "empty.pptx"
    prs.save(str(path))
    from ingestion.pptx_loader import PptxLoader
    chunks = PptxLoader().load(str(path))
    assert chunks == []


def test_pptx_chunk_id_present(tmp_path):
    path = _make_pptx(tmp_path, [
        {"title": "T1", "body": "B1", "notes": "N1"},
        {"title": "T2", "body": "B2"},
    ])
    from ingestion.pptx_loader import PptxLoader
    chunks = PptxLoader().load(path)
    ids = [c["metadata"]["chunk_id"] for c in chunks]
    assert ids == list(range(len(chunks)))


def test_pptx_body_and_notes_share_section_hash(tmp_path):
    path = _make_pptx(tmp_path, [{"title": "T", "body": "Body content.", "notes": "Speaker notes."}])
    from ingestion.pptx_loader import PptxLoader
    chunks = PptxLoader().load(path)
    body_chunk = next(c for c in chunks if not c["metadata"]["is_speaker_notes"])
    notes_chunk = next(c for c in chunks if c["metadata"]["is_speaker_notes"])
    assert body_chunk["metadata"]["section_hash"] == notes_chunk["metadata"]["section_hash"]


def test_pptx_title_excluded_from_body(tmp_path):
    path = _make_pptx(tmp_path, [{"title": "SECRET TITLE", "body": "Actual body content."}])
    from ingestion.pptx_loader import PptxLoader
    chunks = PptxLoader().load(path)
    body_chunk = next(c for c in chunks if not c["metadata"]["is_speaker_notes"])
    assert "SECRET TITLE" not in body_chunk["text"]
    assert "Actual body content." in body_chunk["text"]


def test_pptx_notes_only_no_empty_body_chunk(tmp_path):
    prs = Presentation()
    layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = "Title Only"
    # Leave body placeholder empty
    slide.placeholders[1].text = ""
    slide.notes_slide.notes_text_frame.text = "Only notes here."
    path = tmp_path / "notes_only.pptx"
    prs.save(str(path))
    from ingestion.pptx_loader import PptxLoader
    chunks = PptxLoader().load(str(path))
    body_chunks = [c for c in chunks if not c["metadata"]["is_speaker_notes"]]
    notes_chunks = [c for c in chunks if c["metadata"]["is_speaker_notes"]]
    # No empty body chunk should be emitted
    assert len(body_chunks) == 0
    assert len(notes_chunks) == 1
    assert "Only notes here." in notes_chunks[0]["text"]
