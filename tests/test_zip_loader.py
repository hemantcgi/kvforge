import sys, zipfile
from pathlib import Path
sys.path.insert(0, str(Path(__file__).parent.parent))

from docx import Document
from pptx import Presentation
import openpyxl


def _make_zip(tmp_path, files: dict) -> str:
    """files: {filename: bytes_or_path}"""
    zip_path = tmp_path / "archive.zip"
    with zipfile.ZipFile(zip_path, "w") as zf:
        for name, content in files.items():
            if isinstance(content, (str, Path)):
                zf.write(content, arcname=name)
            else:
                zf.writestr(name, content)
    return str(zip_path)


def test_zip_dispatches_docx(tmp_path):
    doc = Document()
    doc.add_heading("Chapter", level=1)
    doc.add_paragraph("Docx content.")
    docx_path = tmp_path / "doc.docx"
    doc.save(str(docx_path))
    zip_path = _make_zip(tmp_path, {"doc.docx": docx_path})
    from ingestion.zip_loader import ZipLoader
    chunks = ZipLoader().load(zip_path)
    assert any("Docx content." in c["text"] for c in chunks)


def test_zip_dispatches_pptx(tmp_path):
    prs = Presentation()
    layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = "Slide Title"
    slide.placeholders[1].text = "Slide body."
    pptx_path = tmp_path / "deck.pptx"
    prs.save(str(pptx_path))
    zip_path = _make_zip(tmp_path, {"deck.pptx": pptx_path})
    from ingestion.zip_loader import ZipLoader
    chunks = ZipLoader().load(zip_path)
    assert any("Slide body." in c["text"] for c in chunks)


def test_zip_dispatches_xlsx(tmp_path):
    wb = openpyxl.Workbook()
    ws = wb.active
    ws.append(["Product", "Price"])
    ws.append(["Widget", 9.99])
    xlsx_path = tmp_path / "data.xlsx"
    wb.save(str(xlsx_path))
    zip_path = _make_zip(tmp_path, {"data.xlsx": xlsx_path})
    from ingestion.zip_loader import ZipLoader
    chunks = ZipLoader().load(zip_path)
    assert any("Widget" in c["text"] for c in chunks)


def test_zip_skips_unsupported_extensions(tmp_path):
    zip_path = _make_zip(tmp_path, {"readme.txt": b"ignore me"})
    from ingestion.zip_loader import ZipLoader
    chunks = ZipLoader().load(zip_path)
    assert chunks == []


def test_zip_multiple_files(tmp_path):
    doc = Document()
    doc.add_paragraph("From docx.")
    docx_path = tmp_path / "a.docx"
    doc.save(str(docx_path))
    wb = openpyxl.Workbook()
    ws = wb.active
    ws.append(["Col"])
    ws.append(["val"])
    xlsx_path = tmp_path / "b.xlsx"
    wb.save(str(xlsx_path))
    zip_path = _make_zip(tmp_path, {"a.docx": docx_path, "b.xlsx": xlsx_path})
    from ingestion.zip_loader import ZipLoader
    chunks = ZipLoader().load(zip_path)
    texts = " ".join(c["text"] for c in chunks)
    assert "From docx." in texts
    assert "val" in texts


def test_zip_nested_zip_skipped(tmp_path):
    inner_zip = tmp_path / "inner.zip"
    with zipfile.ZipFile(inner_zip, "w") as zf:
        zf.writestr("test.txt", "nested content")
    zip_path = _make_zip(tmp_path, {"inner.zip": inner_zip})
    from ingestion.zip_loader import ZipLoader
    chunks = ZipLoader().load(zip_path)
    assert chunks == []  # nested ZIPs not recursed


def test_zip_chunk_ids_are_sequential(tmp_path):
    doc = Document()
    doc.add_paragraph("Para one.")
    doc.add_paragraph("Para two.")
    docx_path = tmp_path / "doc.docx"
    doc.save(str(docx_path))
    zip_path = _make_zip(tmp_path, {"doc.docx": docx_path})
    from ingestion.zip_loader import ZipLoader
    chunks = ZipLoader().load(zip_path)
    assert len(chunks) >= 1
    # chunk_ids come from the delegated loaders — just verify they exist
    for c in chunks:
        assert "chunk_id" in c["metadata"]


def test_zip_forwards_loader_kwargs(tmp_path):
    # DocxLoader accepts chunk_size kwarg — verify it's forwarded
    doc = Document()
    doc.add_heading("Section", level=1)
    # 200 words — with chunk_size=50, should produce multiple chunks
    doc.add_paragraph(" ".join(f"word{i}" for i in range(200)))
    docx_path = tmp_path / "big.docx"
    doc.save(str(docx_path))
    zip_path = _make_zip(tmp_path, {"big.docx": docx_path})
    from ingestion.zip_loader import ZipLoader
    chunks_small = ZipLoader(chunk_size=50).load(zip_path)
    chunks_large = ZipLoader(chunk_size=1000).load(zip_path)
    # Smaller chunk_size → more chunks
    assert len(chunks_small) > len(chunks_large)


def test_zip_temp_dir_cleaned_up(tmp_path):
    import tempfile
    import unittest.mock as mock
    doc = Document()
    doc.add_paragraph("Some content.")
    docx_path = tmp_path / "doc.docx"
    doc.save(str(docx_path))
    zip_path = _make_zip(tmp_path, {"doc.docx": docx_path})

    # Track temp dirs created
    created_dirs = []
    original_tmpdir = tempfile.TemporaryDirectory

    class TrackingTmpDir:
        def __init__(self):
            self._real = original_tmpdir()
            self.name = self._real.name
            created_dirs.append(self.name)
        def __enter__(self):
            self._real.__enter__()
            return self.name  # Return the string path, not self
        def __exit__(self, *args):
            self._real.__exit__(*args)

    from ingestion.zip_loader import ZipLoader
    import tempfile as tf
    with mock.patch.object(tf, "TemporaryDirectory", TrackingTmpDir):
        ZipLoader().load(zip_path)

    # After load(), all created temp dirs should be gone
    for d in created_dirs:
        assert not Path(d).exists(), f"Temp dir {d} was not cleaned up"
